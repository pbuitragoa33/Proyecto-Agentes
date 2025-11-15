# Generador de reporte


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ==== Paleta de colores / tema ====
COLOR_PRIMARIO = RGBColor(0, 82, 155)      # Azul oscuro
COLOR_ACENTO = RGBColor(0, 160, 150)       # Verde-azulado
COLOR_TEXTO = RGBColor(55, 55, 55)         # Gris oscuro
COLOR_FONDO_CLARO = RGBColor(242, 242, 242)  # Gris muy claro

def estilo_titulo_slide(slide, color=COLOR_PRIMARIO, size=36):
    """Aplica estilo consistente al título de una diapositiva."""
    title_shape = slide.shapes.title
    if not title_shape:
        return

    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = "Calibri"
    font.bold = True
    font.size = Pt(size)
    font.color.rgb = color


def estilo_texto_cuerpo(paragraph, size=18, bold=False, color=COLOR_TEXTO):
    """Aplica estilo básico a un párrafo de cuerpo."""
    font = paragraph.font
    font.name = "Calibri"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color

def estilo_fondo_contenido(slide, color=COLOR_FONDO_CLARO):
    """Aplica un fondo sólido claro a la diapositiva."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def crear_presentacion(nombre_programa: str, datos_snies: dict, datos_agente: dict, output_file: str):
    """
    Crea una presentación de PowerPoint con los resultados del análisis.
    """
    print(f"Generando presentación: {output_file}...")

    prs = Presentation()


    # Diapositiva Título

    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    # Fondo más oscuro para portada
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 80)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Análisis de Oportunidad"
    subtitle.text = nombre_programa

    # Estilo título
    tf_title = title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    font_title = p_title.font
    font_title.name = "Calibri"
    font_title.bold = True
    font_title.size = Pt(44)
    font_title.color.rgb = RGBColor(255, 255, 255)

    # Estilo subtítulo
    tf_sub = subtitle.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    font_sub = p_sub.font
    font_sub.name = "Calibri"
    font_sub.size = Pt(24)
    font_sub.color.rgb = RGBColor(220, 230, 240)

    # Diapositiva: Gráfica Costo vs Matriculados (SNIES) 

    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    estilo_fondo_contenido(slide)
    slide.shapes.title.text = "Análisis SNIES: Costo vs. Matriculados (Colombia)"

    img_path = datos_snies['graficas'].get('costo_vs_matriculados')

    if img_path:

        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width = Inches(8))

    # Diapositiva: Gráfica Programas por Departamento (SNIES)

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    estilo_fondo_contenido(slide)
    slide.shapes.title.text = "Análisis SNIES: Programas por Depto. (Top 10)"
    estilo_titulo_slide(slide)

    img_path = datos_snies['graficas'].get('por_dpto')
    
    if img_path:

        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width = Inches(8))

    # Diapositiva: Gráfica Evolución Estudiantes (SNIES) 

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    estilo_fondo_contenido(slide)
    slide.shapes.title.text = "Análisis SNIES: Evolución de Estudiantes (Procesos)"
    estilo_titulo_slide(slide)

    img_path = datos_snies['graficas'].get('estudiantes_tiempo')

    if img_path:

        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width = Inches(8))

    # Diapositiva: Información Internacional (Agente)

    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    estilo_fondo_contenido(slide)
    slide.shapes.title.text = "Benchmark Internacional (LATAM)"

    textbox = slide.shapes.placeholders[1].text_frame
    textbox.clear() 

    items_latam = [item for item in datos_agente.get('items', []) if item.get('country') not in ['USA', 'España', 'Alemania', 'Colombia']]

    if not items_latam:

        p = textbox.paragraphs[0]
        p.text = "No se encontraron programas en LATAM."

    for item in items_latam:

        p = textbox.add_paragraph()
        p.text = f"{item.get('program_name', 'N/A')} - {item.get('university', 'N/A')} ({item.get('country', 'N/A')})"
        p.font.bold = True
        p.font.size = Pt(18)

        p_cursos = textbox.add_paragraph()
        cursos = ", ".join(item.get('courses_examples', []))
        p_cursos.text = f"  Cursos: {cursos if cursos else 'No disponibles'}"
        p_cursos.level = 1
        p_cursos.font.size = Pt(18)

        p_costo = textbox.add_paragraph()
        p_costo.text = f"  Costo: {item.get('tuition', 'No disponible')}"
        p_costo.level = 1
        p_costo.font.size = Pt(18)


    # Diapositiva: Conclusiones y Tendencias (Agente) 
    
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    estilo_fondo_contenido(slide)
    slide.shapes.title.text = "Análisis de Tendencias y Palabras Clave"
    estilo_titulo_slide(slide)

    textbox = slide.shapes.placeholders[1].text_frame
    textbox.clear()

    insights = datos_agente.get('insights', ["No se generaron insights."])
    
    for insight in insights:

        p = textbox.add_paragraph()
        p.text = f"{insight}"
        p.font.size = Pt(18)

    # Guardar presentación

    prs.save(output_file)
    print(f"Presentación generada exitosamente en: {output_file}")